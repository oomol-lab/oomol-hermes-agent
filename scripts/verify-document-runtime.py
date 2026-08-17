#!/usr/bin/env python3
"""Verify the OOMOL Hermes Agent Office and PDF runtime contract."""

from __future__ import annotations

import importlib
import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


ROOT = Path("/opt/hermes")
SKILLS = Path("/opt/oomol-hermes-agent/curated-skills/productivity")


def require_modules() -> None:
    modules = (
        "markitdown",
        "docx",
        "openpyxl",
        "pptx",
        "defusedxml",
        "reportlab",
        "pypdf",
        "fitz",
        "pyhanko",
        "tzdata",
    )
    for module in modules:
        importlib.import_module(module)


def require_commands() -> None:
    for command in ("soffice", "pandoc", "pdfinfo", "pdftoppm", "qpdf"):
        if shutil.which(command) is None:
            raise RuntimeError(f"required document command is missing: {command}")
    venv_bin = ROOT / ".venv" / "bin"
    if not any((venv_bin / command).is_file() for command in ("pyhanko", "pyhanko-cli")):
        raise RuntimeError("required document command is missing: pyhanko")


def run(*command: str) -> None:
    result = subprocess.run(
        command,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
        timeout=120,
    )
    if result.returncode != 0:
        detail = (result.stderr or result.stdout).strip()
        raise RuntimeError(f"document smoke command failed: {command[0]}: {detail}")


def verify_office(workspace: Path) -> None:
    from docx import Document
    from markitdown import MarkItDown
    from openpyxl import Workbook, load_workbook
    from pptx import Presentation

    docx_path = workspace / "office-smoke.docx"
    document = Document()
    document.add_heading("OOMOL 文档验证", level=1)
    document.add_paragraph("Office runtime ready")
    document.save(docx_path)
    reopened = Document(docx_path)
    assert "Office runtime ready" in "\n".join(
        paragraph.text for paragraph in reopened.paragraphs
    )
    extracted = MarkItDown().convert(docx_path).text_content
    assert "Office runtime ready" in extracted

    xlsx_path = workspace / "office-smoke.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "OOMOL Excel 验证"
    workbook.save(xlsx_path)
    assert load_workbook(xlsx_path).active["A1"].value == "OOMOL Excel 验证"

    pptx_path = workspace / "office-smoke.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "OOMOL PowerPoint 验证"
    slide.placeholders[1].text = "Basic PPTX runtime ready"
    presentation.save(pptx_path)
    assert Presentation(pptx_path).slides[0].shapes.title.text == "OOMOL PowerPoint 验证"

    markdown_path = workspace / "office-smoke.md"
    markdown_path.write_text("# Pandoc smoke\n", encoding="utf-8")
    run("pandoc", str(markdown_path), "-o", str(workspace / "pandoc-smoke.docx"))

    profile = workspace / "libreoffice-profile"
    profile.mkdir()
    output = workspace / "libreoffice-output"
    output.mkdir()
    run(
        "soffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--norestore",
        f"-env:UserInstallation={profile.as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output),
        str(docx_path),
    )
    converted = output / "office-smoke.pdf"
    if not converted.is_file() or converted.stat().st_size == 0:
        raise RuntimeError("LibreOffice did not produce a non-empty PDF")
    run("qpdf", "--check", str(converted))


def verify_pdf(workspace: Path) -> None:
    from pypdf import PdfReader

    pdf_skill = SKILLS / "pdf-files"
    scripts = pdf_skill / "scripts"
    required_scripts = (
        "create_labeled_acroform.py",
        "inspect_pdf_widgets.py",
        "inspect_pdf_form_labels.py",
        "fill_acroform_static.py",
        "verify_pdf_delivery.py",
        "sign_pdf_test_cert.py",
    )
    for name in required_scripts:
        if not (scripts / name).is_file():
            raise RuntimeError(f"bundled PDF helper is missing: {name}")

    blank_pdf = workspace / "form.pdf"
    spec = workspace / "form-spec.json"
    spec.write_text(
        json.dumps(
            {
                "output": str(blank_pdf),
                "title": "OOMOL PDF 表单验证",
                "field_rows": [
                    [{"name": "contact", "label": "联系人"}],
                    [
                        {
                            "name": "note",
                            "label": "备注",
                            "kind": "multiline",
                            "height": 56,
                        }
                    ],
                ],
                "footer": "Hermes document runtime smoke",
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    run(sys.executable, str(scripts / "create_labeled_acroform.py"), str(spec))
    run(sys.executable, str(scripts / "inspect_pdf_widgets.py"), str(blank_pdf))
    run(
        sys.executable,
        str(scripts / "inspect_pdf_form_labels.py"),
        str(blank_pdf),
        "--labels",
        json.dumps({"contact": "联系人", "note": "备注"}, ensure_ascii=False),
    )

    filled_pdf = workspace / "filled.pdf"
    values = json.dumps(
        {"contact": "刘志强", "note": "静态 PDF 验证通过"},
        ensure_ascii=False,
    )
    run(
        sys.executable,
        str(scripts / "fill_acroform_static.py"),
        str(blank_pdf),
        str(filled_pdf),
        "--values",
        values,
    )
    run(
        sys.executable,
        str(scripts / "verify_pdf_delivery.py"),
        str(filled_pdf),
        "--contains",
        "刘志强",
        "--no-duplicate-nearby",
        "刘志强",
    )

    signed_pdf = workspace / "signed.pdf"
    run(
        sys.executable,
        str(scripts / "sign_pdf_test_cert.py"),
        str(filled_pdf),
        str(signed_pdf),
        "--field-name",
        "RuntimeSmokeSignature",
    )
    fields = PdfReader(signed_pdf).get_fields() or {}
    signature = fields.get("RuntimeSmokeSignature")
    if not signature or signature.get("/FT") != "/Sig":
        raise RuntimeError("test-signed PDF is missing its signature field")


def main() -> int:
    require_modules()
    require_commands()
    for skill_name in ("office-files", "pdf-files"):
        if not (SKILLS / skill_name / "SKILL.md").is_file():
            raise RuntimeError(f"bundled document skill is missing: {skill_name}")

    with tempfile.TemporaryDirectory(prefix="oomol-document-smoke-") as temp_dir:
        workspace = Path(temp_dir)
        verify_office(workspace)
        verify_pdf(workspace)

    print("OOMOL document runtime: verified")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
